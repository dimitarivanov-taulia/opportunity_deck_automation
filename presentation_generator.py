import json
import os
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
import time
from dotenv import load_dotenv
from bing_image_downloader import downloader
import certifi
import requests
from datetime import datetime

load_dotenv()
# Load the issues information from the JSON file
with open('issues_info.json', 'r') as json_file:
    issues_info = json.load(json_file)

# Assuming we're working with the first issue in the list for now
first_issue = issues_info[0]
GLOBAL_BUYER_NAME = first_issue["Buyer Company Name"]
FILTER_UPDATES = first_issue

# Load the dashboard dictionary from environment variable
DASHBOARD_DICT = json.loads(os.getenv('DASHBOARD_DICT'))

# Set the SSL certificate file environment variable
os.environ['SSL_CERT_FILE'] = certifi.where()

# Set up your OpenAI API key and Bing Image Search API key
openai.api_key = os.getenv('OPENAI_API_KEY')
bing_image_api_key = os.getenv('BING_IMAGE_API_KEY')
scale_serp_api_key = os.getenv('SCALE_SERP_API_KEY')
buyer_name = GLOBAL_BUYER_NAME
images_folder = 'images'  # Folder where logos and other images are stored
data_folder = 'data'  # Folder where the data files are stored
open_ai_assistant_name = 'Presentation Generator'
last_message_timestamp = 0

assistant_subtitle_questions = {
    'Peer DPO': f"Based on the provided csv data how many days is {buyer_name} behind or above the peer group average? Give short sentence only with the answer.",
    'Industry Term Comparison': f"Based on the provided csv data how many days is {buyer_name} behind or above industry norms across top 10 industries? Give short sentence only with the answer.",
    'Working Capital Released': f"Based on the provided csv data what working capital can {buyer_name} unlock? Give short sentence only with the answer and format number like MM."
}

def get_assistant_id_by_name(client, assistant_name):
    list_assistants = client.beta.assistants.list()
    for assistant in list_assistants:
        if assistant.name == assistant_name:
            return assistant.id
    raise Exception(f"No Assistant with name '{assistant_name}' found on OpenAI!")

def start_new_conversation(client):
    new_thread = client.beta.threads.create()
    return new_thread

def send_message(client, conversation, assistant_id, content):
    print(f"Sending message to assistant:\n{content}")
    client.beta.threads.messages.create(
        thread_id=conversation.id,
        role="user",
        content=content
    )
    client.beta.threads.runs.create_and_poll(
        thread_id=conversation.id,
        assistant_id=assistant_id
    )

def has_message_completed(message):
    if 'content' not in message.dict():
        return False
    if len(message.dict()['content']) < 1:
        return False
    if 'text' not in message.dict()['content'][0]:
        return False
    if 'value' not in message.dict()['content'][0]['text']:
        return False
    return True

def get_assistant_response(client, conversation, last_message_timestamp):
    try:
        messages = client.beta.threads.messages.list(thread_id=conversation.id, limit=20)
        if messages is None:
            return None

        for message in messages:
            if message.role == 'assistant' and has_message_completed(message):
                if message.created_at > last_message_timestamp:
                    last_message_timestamp = message.created_at
                    return message
        return None
    except Exception as e:
        print(f"Unable to retrieve assistant response: {e}")
        return None

def get_summary_from_assistant(client, assistant_name, data_folder, buyer_name):
    assistant_id = get_assistant_id_by_name(client, assistant_name)
    conversation = start_new_conversation(client)

    content = f"Buyer's name: {buyer_name}\nData:\n"
    for filename in os.listdir(data_folder):
        if filename.endswith('_with_filters.csv'):
            with open(os.path.join(data_folder, filename), 'r') as file:
                csv_content = file.read()
                content += f"\nCSV file: {filename}\n{csv_content}\n"

    send_message(client, conversation, assistant_id, content)

    response = None
    while response is None:
        response = get_assistant_response(client, conversation, last_message_timestamp)
        time.sleep(2)

    return response.dict()['content'][0]['text']['value']

def get_subtitle_from_assistant(client, assistant_name, element_name, data_file_path):
    assistant_id = get_assistant_id_by_name(client, assistant_name)
    conversation = start_new_conversation(client)

    with open(data_file_path, 'r') as file:
        csv_content = file.read()

    content = f"Buyer's name: {buyer_name}\nCSV file content:\n{csv_content}\n{assistant_subtitle_questions[element_name]}"

    send_message(client, conversation, assistant_id, content)

    response = None
    while response is None:
        response = get_assistant_response(client, conversation, last_message_timestamp)
        time.sleep(2)

    return response.dict()['content'][0]['text']['value']

def download_image(query, download_folder, limit=1):
    downloader.download(query, limit=limit, output_dir=download_folder, adult_filter_off=True, force_replace=False, timeout=60)
    downloaded_images = os.listdir(os.path.join(download_folder, query))
    if downloaded_images:
        return os.path.join(download_folder, query, downloaded_images[0])
    return None

def search_text(query, api_key):
    url = "https://api.scaleserp.com/search"
    params = {
        "api_key": api_key,
        "q": query,
        "location": "United States",
        "hl": "en",
        "gl": "us"
    }
    response = requests.get(url, params=params)
    results = response.json().get('organic_results', [])
    if results:
        for result in results:
            if 'snippet' in result:
                return result['snippet']
    return None

def extract_industry_from_snippet(snippet):
    keyword = 'industry'
    if keyword in snippet:
        parts = snippet.split(keyword)[0].strip().split()
        return ' '.join(parts)
    return "General"

def download_buyer_logo(buyer_name, download_folder):
    return download_image(f"{buyer_name} logo", download_folder)

def download_industry_background(industry, download_folder):
    return download_image(f"{industry} background", download_folder)

def get_buyer_industry(buyer_name, api_key):
    snippet = search_text(f"{buyer_name} industry", api_key)
    if snippet:
        return extract_industry_from_snippet(snippet)
    return "General"

def add_title_slide(prs, logo_path, background_path, buyer_logo_path=None):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(0)
    top = Inches(0)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide.shapes.add_picture(background_path, left, top, width=slide_width, height=slide_height)

    if buyer_logo_path:
        buyer_logo_left = Inches(1)
        buyer_logo_top = Inches(3.5)
        buyer_logo_height = Inches(2)
        slide.shapes.add_picture(buyer_logo_path, buyer_logo_left, buyer_logo_top, height=buyer_logo_height)

    logo_left = Inches(7.75)
    logo_top = Inches(0.5)
    logo_height = Inches(0.4)
    slide.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)
    date_time = datetime.now()
    title_placeholder = "Terms Extension Opportunity Assessment"
    subtitle_placeholder = f"{date_time.strftime('%B %Y')}"

    text_left = Inches(1)
    text_top = Inches(5.5)
    text_width = Inches(7)
    text_height = Inches(0.5)
    title_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    title_frame = title_box.text_frame
    title_frame.text = title_placeholder
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    text_top = Inches(6)
    subtitle_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle_placeholder
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

def add_footer(prs, slide):
    slide_width = prs.slide_width

    footer_text = "© 2023 Taulia LLC. Confidential information. Do not distribute."
    footer_box = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.22), slide_width / 2 - Inches(0.5), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = footer_text
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    bar_left = slide_width / 2
    bar_top = prs.slide_height - Inches(0.15)
    bar_width = slide_width / 2
    bar_height = Inches(0.15)
    bar = slide.shapes.add_shape(1, bar_left, bar_top, bar_width, bar_height)
    fill = bar.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 102, 0)
    bar.line.fill.background()

def add_summary_slide(prs, summary):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    text_left = Inches(1)
    text_top = Inches(0.17)
    text_width = Inches(3)
    text_height = Inches(0.5)
    title_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    title_frame = title_box.text_frame
    title_frame.text = "Situation Summary"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True

    text_left = Inches(1)
    text_top = Inches(1)
    text_width = prs.slide_width - Inches(2)
    text_height = prs.slide_height - Inches(2)
    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = text_box.text_frame
    text_frame.text = summary
    text_frame.word_wrap = True

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(11)
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    add_footer(prs, slide)

def add_content_slide(prs, title, img_path, logo_path, subtitle=None):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.title
    title_box.text = title

    for paragraph in title_box.text_frame.paragraphs:
        paragraph.font.size = Pt(20)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.word_wrap = True
        subtitle_frame.paragraphs[0].font.size = Pt(14)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 102, 0)
        subtitle_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    left = Inches(1)
    top = Inches(2)
    height = Inches(4.5)
    slide.shapes.add_picture(img_path, left, top, height=height)

    logo_height = Inches(0.3)
    slide_width = prs.slide_width
    logo_left = slide_width - Inches(2)
    logo_top = Inches(0.4)
    slide.shapes.add_picture(logo_path, logo_left, logo_top, height=logo_height)

    add_footer(prs, slide)

def create_ppt(images_folder, output_ppt, logo_path, assistant_name, buyer_name):
    prs = Presentation()

    print("Finding buyer industry...")
    buyer_industry = get_buyer_industry(buyer_name, scale_serp_api_key)
    print(f"Buyer industry found: {buyer_industry}")

    print("Downloading industry background image...")
    background_path = download_industry_background(buyer_industry, images_folder)
    if not background_path:
        print("Error: Unable to download background image. Using default background.")
        background_path = os.path.join(images_folder, 'background.jpg')
    print(f"Background image set to {background_path}")

    print("Downloading buyer logo...")
    buyer_logo_path = download_buyer_logo(buyer_name, images_folder)
    if buyer_logo_path:
        print(f"Buyer logo downloaded to {buyer_logo_path}")
    else:
        print("Buyer logo not found. Skipping buyer logo.")

    add_title_slide(prs, logo_path, background_path, buyer_logo_path)

    print("Getting summary from assistant...")
    summary = get_summary_from_assistant(openai, assistant_name, data_folder, buyer_name)
    print("Summary obtained from assistant.")

    add_summary_slide(prs, summary)

    for filename in os.listdir(images_folder):
        if filename.endswith(('.png', '.jpg', '.jpeg')) and not filename.startswith('taulia-logo') and not filename.startswith('background'):
            slide_title = os.path.splitext(filename)[0]
            img_path = os.path.join(images_folder, filename)
            subtitle = None
            if slide_title in assistant_subtitle_questions:
                data_file_path = os.path.join(data_folder, f"{slide_title}_with_filters.csv")
                if os.path.exists(data_file_path):
                    subtitle = get_subtitle_from_assistant(openai, assistant_name, slide_title, data_file_path)
            add_content_slide(prs, slide_title, img_path, logo_path, subtitle)

    try:
        prs.save(output_ppt)
        print(f"Presentation saved as {output_ppt}")
    except Exception as e:
        print(f"Error saving presentation: {e}")

images_folder = 'images'
logo_path = os.path.join(images_folder, 'taulia-logo.png')
data_folder = 'data'  # Folder where the data files are stored
output_ppt = f'presentations/opportunity_deck_{buyer_name}.pptx'
assistant_name = 'Presentation Generator'

create_ppt(images_folder, output_ppt, logo_path, assistant_name, buyer_name)
